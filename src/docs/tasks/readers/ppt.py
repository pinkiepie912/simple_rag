from typing import List, cast

from llama_index.core.schema import Document
from pptx import Presentation
from pptx.shapes.autoshape import Shape


class PptReader:
    def load_data(self, file_path: str):
        prs = Presentation(file_path)
        documents: List[Document] = []

        for idx, slide in enumerate(prs.slides):
            lines = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                text_shape = cast(Shape, shape)
                for paragraph in text_shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in paragraph.runs).strip()
                    if text:
                        lines.append(text)

            if lines:
                full_text = '\n'.join(lines)
                documents.append(Document(text=full_text, metadata={"slide_index": idx}))

        return documents
